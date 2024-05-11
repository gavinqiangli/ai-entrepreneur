from pptx import Presentation
import streamlit as st
import json


def generate_ppt():
    # Your slide_content JSON data
    slide_content_json = """
    [
    {
      "slide_title": "Problem",
      "main_bullet": "Complex and Time-Consuming Process",
      "child_bullet": "High costs with patent attorneys, Long wait times, High error rates in manual drafting"
    },
    {
      "slide_title": "Solution",
      "main_bullet": "AI-Driven Patent Drafting",
      "child_bullet": "Automates drafting, Provides real-time suggestions, Ensures compliance with laws"
    },
    {
      "slide_title": "Market Opportunity",
      "main_bullet": "Expansive and Growing Market",
      "child_bullet": "Global IP software market to reach $5B by 2026, Targeting inventors and corporations"
    },
    {
      "slide_title": "Product Demo",
      "main_bullet": "Live Demonstration",
      "child_bullet": "Showcase of drafting tool, Real-time assistance, Error correction features"
    },
    {
      "slide_title": "Customer Testimonials",
      "main_bullet": "Positive Feedback",
      "child_bullet": "Ease of use, Cost savings, Improved speed"
    },
    {
      "slide_title": "Team",
      "main_bullet": "Experienced and Diverse",
      "child_bullet": "AI researchers, Patent attorneys, Software developers"
    },
    {
      "slide_title": "Revenue and Growth",
      "main_bullet": "Strong Business Model",
      "child_bullet": "Subscription-based, Tiered pricing, 5% market capture in 3 years"
    },
    {
      "slide_title": "Competitive Landscape",
      "main_bullet": "Unique Selling Proposition",
      "child_bullet": "Faster, cheaper, more accurate than competitors, AI integration"
    },
    {
      "slide_title": "Conclusion",
      "main_bullet": "Investment Opportunity",
      "child_bullet": "Scale technology, Expand market reach, Democratize patent services"
    }
    ]
    """

    # Load JSON data
    slide_content = json.loads(slide_content_json)

    # Hello World! example
    # ../_images/hello-world.png
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

    # Bullet slide example
    # ../_images/bullet-slide.png
    # Iterate through each slide and print its content
    for slide in slide_content:
        slide_title = slide["slide_title"]
        main_bullet = slide["main_bullet"]
        child_bullet = slide.get('child_bullet', None)  # Default value is None if key doesn't exist
        grandchild_bullet = slide.get('grandchild_bullet', None)  # Default value is None if key doesn't exist

        print("Slide Title:", slide_title)
        print("Main Bullet:", main_bullet)
        print("Child Bullet:", child_bullet)
        print("Grandchild Bullet:", grandchild_bullet)
        print("\n")

        bullet_slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = slide_title

        tf = body_shape.text_frame
        tf.text = main_bullet

        # Check if child_bullet is a list
        if child_bullet is not None:
            if isinstance(child_bullet, list):
                for bullet in child_bullet:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 1
            else:
                p = tf.add_paragraph()
                p.text = child_bullet
                p.level = 1

        # Check if grandchild_bullet is a list
        if grandchild_bullet is not None:
            if isinstance(grandchild_bullet, list):
                for bullet in grandchild_bullet:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 2
            else:
                p = tf.add_paragraph()
                p.text = grandchild_bullet
                p.level = 2

    return prs


# save locally
# prs = generate_file()
# prs.save('test.pptx')


# save on streamlit cloud (download file to your web browser)
@st.cache_resource(ttl="1h")
def save_file():
    temp_filepath = "test.pptx"
    prs = generate_ppt()
    prs.save(temp_filepath)
    return temp_filepath


temp_filepath = save_file()
with open(temp_filepath, "rb") as template_file:
    template_byte = template_file.read()

st.download_button(
    label="Click to download your test PPT",
    data=template_byte,
    file_name="test.pptx",
    mime="application/octet-stream",
)
